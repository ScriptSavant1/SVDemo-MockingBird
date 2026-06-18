"""PowerPoint report renderer using python-pptx.

Produces a 6-slide management deck:
  1. Title slide (project name, stub name, date)
  2. Executive Summary (4 KPI boxes)
  3. Performance Overview (peak/avg TPS, latency)
  4. Error Analysis (error rate, total errors)
  5. Top 10 Time-Series data points (table)
  6. Recommendations + Next Steps
"""
from __future__ import annotations

import io
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import ReportData


def _rgb(hex_colour: str) -> RGBColor:
    h = hex_colour.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    colour: str = "#222222",
    align=PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(colour)


def _add_kpi_box(
    slide, label: str, value: str, left: float, top: float, primary: str
) -> None:
    """Coloured 1.8 × 1.2 inch KPI block."""
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(1.8), Inches(1.2),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(primary)
    bg.line.fill.background()

    # Value text
    tf_val = bg.text_frame
    tf_val.word_wrap = False
    p_val = tf_val.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    run_val = p_val.add_run()
    run_val.text = value
    run_val.font.size = Pt(24)
    run_val.font.bold = True
    run_val.font.color.rgb = _rgb("#FFFFFF")

    # Label text box below
    _add_text_box(
        slide, label,
        left, top + 1.25, 1.8, 0.4,
        font_size=10, bold=False, colour="#555555",
        align=PP_ALIGN.CENTER,
    )


def render_ppt(data: ReportData, primary: str, secondary: str, company: str) -> bytes:
    """Return PowerPoint bytes (widescreen 13.33 × 7.5 in)."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(primary)

    _add_text_box(s1, company, 0.5, 0.3, 12, 0.6, 14, False, "#FFFFFF")
    _add_text_box(s1, "Stub Engine Performance Report", 0.5, 1.2, 12, 1.0, 36, True, "#FFFFFF", PP_ALIGN.CENTER)
    _add_text_box(s1, f"Project: {data.project_name}  |  Stub: {data.stub_name}", 0.5, 2.5, 12, 0.6, 18, False, "#FFFFFF", PP_ALIGN.CENTER)
    _add_text_box(s1, f"Generated: {data.generated_at.strftime('%d %B %Y')}", 0.5, 3.2, 12, 0.5, 14, False, "#CCDDEE", PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _add_text_box(s2, "Executive Summary", 0.5, 0.2, 12, 0.7, 28, True, primary)

    kpis = [
        ("Peak TPS", f"{data.peak_tps:,.0f}"),
        ("Avg TPS", f"{data.avg_tps:,.0f}"),
        ("Avg Latency", f"{data.avg_latency_ms:.1f}ms"),
        ("Error Rate", f"{data.error_rate_pct:.3f}%"),
    ]
    for i, (label, val) in enumerate(kpis):
        _add_kpi_box(s2, label, val, 1.0 + i * 2.8, 1.2, primary)

    _add_text_box(s2, f"Environment: {data.environment}  |  Period: last {data.report_period_hours}h  |  Total requests: {data.total_requests:,}", 0.5, 5.8, 12, 0.5, 12, False, "#555555")

    # ── Slide 3: Performance ──────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    _add_text_box(s3, "Performance Overview", 0.5, 0.2, 12, 0.7, 28, True, primary)
    rows = [
        ("Metric", "Value"),
        ("Peak TPS", f"{data.peak_tps:,.1f}"),
        ("Average TPS", f"{data.avg_tps:,.1f}"),
        ("Avg Latency (ms)", f"{data.avg_latency_ms:.2f}"),
        ("p95 Latency (ms)", f"{data.p95_latency_ms:.2f}"),
        ("Stub URL", data.stub_url),
    ]
    _add_perf_table(s3, rows, primary, secondary)

    # ── Slide 4: Error Analysis ───────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    _add_text_box(s4, "Error Analysis", 0.5, 0.2, 12, 0.7, 28, True, primary)
    error_rows = [
        ("Metric", "Value"),
        ("Error Rate (%)", f"{data.error_rate_pct:.4f}"),
        ("Total Errors (est.)", f"{int(data.total_requests * data.error_rate_pct / 100):,}"),
        ("Total Requests", f"{data.total_requests:,}"),
    ]
    _add_perf_table(s4, error_rows, primary, secondary)

    # ── Slide 5: Time-Series Sample ───────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    _add_text_box(s5, f"Time-Series Data (first {min(10, len(data.points))} of {len(data.points)} observations)", 0.5, 0.2, 12, 0.7, 22, True, primary)
    ts_rows = [("Time (UTC)", "TPS", "Avg Latency (ms)", "Error Rate (%)")]
    for p in data.points[:10]:
        ts_rows.append((
            p.time.strftime("%Y-%m-%d %H:%M"),
            f"{p.tps:,.1f}",
            f"{p.latency_avg_ms:.2f}",
            f"{p.error_rate * 100:.4f}",
        ))
    _add_wide_table(s5, ts_rows, primary, secondary)

    # ── Slide 6: Recommendations ──────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    _add_text_box(s6, "Recommendations & Next Steps", 0.5, 0.2, 12, 0.7, 28, True, primary)
    bullets = [
        f"• Current peak TPS {data.peak_tps:,.0f} — target is 10,000 TPS per stub.",
        "• Monitor p95 latency trend; investigate if > 100ms sustained.",
        f"• Error rate {data.error_rate_pct:.3f}% — escalate if > 0.1% in production.",
        "• Stub stubs are preserved in PostgreSQL + S3 — redeploy in < 4 minutes.",
        "• Download full Excel report for raw data analysis.",
    ]
    for i, bullet in enumerate(bullets):
        _add_text_box(s6, bullet, 0.8, 1.2 + i * 0.9, 11.5, 0.8, 16, False, "#333333")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_perf_table(slide, rows: list[tuple], primary: str, secondary: str) -> None:
    """Add a 2-column table to the slide starting at (1.5, 1.2)."""
    if not rows:
        return
    tbl = slide.shapes.add_table(len(rows), 2, Inches(1.5), Inches(1.2), Inches(7), Inches(0.5 * len(rows))).table
    for ri, (c1, c2) in enumerate(rows):
        _style_cell(tbl.cell(ri, 0), c1, primary if ri == 0 else None, bold=(ri == 0))
        _style_cell(tbl.cell(ri, 1), c2, primary if ri == 0 else None, bold=(ri == 0))


def _add_wide_table(slide, rows: list[tuple], primary: str, secondary: str) -> None:
    if not rows:
        return
    tbl = slide.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.2), Inches(12), Inches(0.45 * len(rows))).table
    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            _style_cell(tbl.cell(ri, ci), val, primary if ri == 0 else None, bold=(ri == 0))


def _style_cell(cell, text: str, bg_hex: str | None, bold: bool = False) -> None:
    cell.text = text
    tf = cell.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(11)
        run.font.bold = bold
        if bg_hex:
            run.font.color.rgb = _rgb("#FFFFFF")
    if bg_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg_hex)
